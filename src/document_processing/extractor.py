"""
File content extractor supporting multiple formats.

Each extractor receives raw bytes and returns plain text.
"""

import csv
import io
import json
import logging

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".rst",
    ".pdf",
    ".docx",
    ".xlsx", ".xls",
    ".pptx",
    ".csv",
    ".json",
    ".yaml", ".yml",
    ".html", ".htm",
}


def extract_text(content: bytes, extension: str) -> str:
    """
    Extract plain text from file bytes based on extension.

    Raises:
        ValueError: If the format is unsupported or content cannot be parsed.
    """
    ext = extension.lower()

    if ext in (".txt", ".md", ".rst"):
        return _extract_plaintext(content)
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(content)
    if ext == ".pptx":
        return _extract_pptx(content)
    if ext == ".csv":
        return _extract_csv(content)
    if ext == ".json":
        return _extract_json(content)
    if ext in (".yaml", ".yml"):
        return _extract_yaml(content)
    if ext in (".html", ".htm"):
        return _extract_html(content)

    raise ValueError(f"Unsupported file format: {ext}")


# ---------------------------------------------------------------------------
# Format handlers
# ---------------------------------------------------------------------------

def _extract_plaintext(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


def _extract_pdf(content: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text.strip()}")
    if not pages:
        raise ValueError("PDF contains no extractable text (may be scanned/image-only)")
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(content))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Include tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    if not parts:
        raise ValueError("DOCX contains no extractable text")
    return "\n\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    if not sheets:
        raise ValueError("Spreadsheet contains no data")
    return "\n\n".join(sheets)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    if not slides:
        raise ValueError("Presentation contains no extractable text")
    return "\n\n".join(slides)


def _extract_csv(content: bytes) -> str:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1")
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(cell.strip() for cell in row) for row in reader if any(c.strip() for c in row)]
    if not rows:
        raise ValueError("CSV contains no data")
    return "\n".join(rows)


def _extract_json(content: bytes) -> str:
    try:
        data = json.loads(content.decode("utf-8"))
    except (json.JSONDecodeError, UnicodeDecodeError) as e:
        raise ValueError(f"Invalid JSON: {e}")
    return json.dumps(data, indent=2, ensure_ascii=False)


def _extract_yaml(content: bytes) -> str:
    import yaml
    try:
        data = yaml.safe_load(content.decode("utf-8"))
    except Exception as e:
        raise ValueError(f"Invalid YAML: {e}")
    return yaml.dump(data, allow_unicode=True, default_flow_style=False)


def _extract_html(content: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style", "head", "meta", "link"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        raise ValueError("HTML contains no extractable text")
    return "\n".join(lines)
